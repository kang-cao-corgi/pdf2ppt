import os
import platform
from typing import Sequence, List, Optional
from PIL import Image
from pdf2image import convert_from_path  # pyright: ignore reportUnknownVariableType
from pptx import Presentation  # pyright: ignore reportUnknownVariableType
from io import BytesIO
from pathlib import PurePath


DEBUG = 0
PIXEL_TO_EMU_RATIO = 3_000
DPI = 96


def pixels_to_emu(pixels: int):
      return pixels * PIXEL_TO_EMU_RATIO  # TODO: finer control


def save_images(
    image_list: Sequence[Image.Image],
    output_folder: str ='output_images',
    file_prefix: str ='image'
):
    """
    Save a list of PIL image objects to individual files.

    Parameters:
    - image_list: List of PIL image objects.
    - output_folder: Output folder path where images will be saved.
    - file_prefix: Prefix for the saved image filenames.
    """
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for i, img in enumerate(image_list):
        filename = f"{file_prefix}_{i + 1}.png"
        filepath = os.path.join(output_folder, filename)
        img.save(filepath)


def load_images(
    folder_path: str = 'output_images',
    file_prefix: str = 'image'
) -> Sequence[Image.Image]:
    """
    Load a list of PIL image objects from individual files.

    Parameters:
    - folder_path: Folder path where images are saved.
    - file_prefix: Prefix used for saved image filenames.

    Returns:
    - List of PIL image objects.
    """
    image_list: List[Image.Image] = []

    for filename in sorted(os.listdir(folder_path)):
        if filename.startswith(file_prefix) and filename.endswith('.png'):
            filepath = os.path.join(folder_path, filename)
            img = Image.open(filepath)
            image_list.append(img)

    return image_list


def get_poppler_path() -> Optional[PurePath]:
    if platform.system().lower() == "windows" and (not DEBUG):
        return PurePath(__file__).parent.parent / "poppler/bin"  # XXX: list your own path
    return None  # use default in mac / debug


def convert(
    input_file_path: str,
    output_path: str,
):

	print("Converting file: " + input_file_path)

	# Prep presentation
	prs = Presentation()  # pyright: ignore reportUnknownVariableType
	blank_slide_layout = prs.slide_layouts[6]  # pyright: ignore reportUnknownVariableType

	# Create working folder
	base_name = input_file_path.split("/")[-1].split(".pdf")[0]

	# Convert PDF to list of images
	print("Starting conversion...")

	slideimgs = convert_from_path(input_file_path, DPI, fmt='ppm', thread_count=3, poppler_path=get_poppler_path())  # pyright: ignore reportArgumentType

	print("...complete.")

	# Loop over slides
	for i, slideimg in enumerate(slideimgs):
		if i % 10 == 0:
			print("Saving slide: " + str(i))

		imagefile = BytesIO()
		slideimg.save(imagefile, format='tiff')
		# imagedata = imagefile.getvalue()
		imagefile.seek(0)
		width, height = slideimg.size

		# Set slide dimensions
		height_emu = pixels_to_emu(height)
		width_emu = pixels_to_emu(width)

		prs.slide_height = height_emu
		prs.slide_width = width_emu

		# Add slide
		slide = prs.slides.add_slide(blank_slide_layout)  # pyright: ignore reportUnknownMemberType
		slide.shapes.add_picture(imagefile, 0, 0, width=width_emu, height=height_emu)  # pyright: ignore reportUnknownMemberType

	# Save Powerpoint
	print("Saving file: " + base_name + ".pptx")
	prs.save(f'{output_path}/{base_name}.pptx')  # pyright: ignore reportUnknownMemberType
	print("Conversion complete. :)")
